import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Initialize Presentation
    prs = Presentation()
    
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_bg = RGBColor(13, 13, 13)         # Pure Dark Slate/Black
    c_card = RGBColor(30, 30, 30)       # Dark Card Gray
    c_red = RGBColor(225, 6, 0)         # F1 Racing Red
    c_cyan = RGBColor(102, 252, 241)    # Telemetry Cyan/Teal
    c_white = RGBColor(255, 255, 255)   # Bright White
    c_gray = RGBColor(180, 180, 180)    # Muted Gray
    
    c_green = RGBColor(46, 204, 113)    # Calm Green
    c_orange = RGBColor(243, 156, 18)   # Tired Orange
    
    blank_slide_layout = prs.slide_layouts[6] # Blank layout
    
    def set_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg

    def add_header(slide, title_text, category_text="PITPULSE HACKATHON"):
        # Header background banner line
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = c_red
        shape.line.fill.background()
        
        # Category indicator
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_red
        
        # Main Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.name = "Arial"
        p_title.font.size = Pt(28)
        p_title.font.bold = True
        p_title.font.color.rgb = c_white

    def create_card(slide, left, top, width, height, bg_color):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = RGBColor(60, 60, 60)
        shape.line.width = Pt(1)
        return shape

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide1)
    
    # Large Decorative Red Band on left
    stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = c_red
    stripe.line.fill.background()
    
    # Center-ish Container for title
    title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10), Inches(3.0))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "PITPULSE"
    p1.font.name = "Arial"
    p1.font.size = Pt(64)
    p1.font.bold = True
    p1.font.color.rgb = c_white
    
    p2 = tf1.add_paragraph()
    p2.text = "F1 Driver Acoustic Sentiment & Telemetry Correlation"
    p2.font.name = "Arial"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = c_red
    p2.space_before = Pt(10)
    
    p3 = tf1.add_paragraph()
    p3.text = "Detecting stress, fatigue, and pace loss in real time to power actionable race-engineering decisions."
    p3.font.name = "Arial"
    p3.font.size = Pt(14)
    p3.font.color.rgb = c_gray
    p3.space_before = Pt(14)
    
    p4 = tf1.add_paragraph()
    p4.text = "HACKATHON PRESENTATION  |  STAGE PITCH"
    p4.font.name = "Arial"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = c_cyan
    p4.space_before = Pt(40)

    # ==========================================
    # SLIDE 2: THE PROBLEM (Grid Layout)
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide2)
    add_header(slide2, "The Problem: The Driver's Invisible Struggle")
    
    # Card 1: Extreme Driver Load
    create_card(slide2, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), c_card)
    tb_c1 = slide2.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_left = tf_c1.margin_right = tf_c1.margin_top = tf_c1.margin_bottom = 0
    
    p = tf_c1.paragraphs[0]
    p.text = "HIGH-STRESS ENVIRONMENT"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(15)
    
    bullets_c1 = [
        ("Extreme Physics:", " Drivers experience up to 5G lateral forces, heart rates averaging 170+ BPM, and dehydrate up to 4kg per race."),
        ("Cognitive Overload:", " Deciding strategy at 300+ km/h leaves zero margin for cognitive fatigue or emotional stress."),
        ("Micro-Decisions:", " A fraction of a second of delay in reaction times due to mental strain leads directly to lost pace or catastrophic crashes.")
    ]
    for title, text in bullets_c1:
        p = tf_c1.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = c_white
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(13)
        r2.font.color.rgb = c_gray
        
    # Card 2: Technical Blind Spot
    create_card(slide2, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8), c_card)
    tb_c2 = slide2.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_left = tf_c2.margin_right = tf_c2.margin_top = tf_c2.margin_bottom = 0
    
    p = tf_c2.paragraphs[0]
    p.text = "THE PIT WALL BLIND SPOT"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(15)
    
    bullets_c2 = [
        ("Telemetry Isolation:", " Race engineers monitor vehicle telemetry (tire pressure, engine temp, sector times) but have no real-time data on driver stress."),
        ("Subjective Comms:", " Radio communication analysis is completely subjective, relying on the engineer's 'gut feeling' under high-pressure scenarios."),
        ("Acoustic Noise:", " Heavy background engine rumble (V6 Hybrid) and wind noise mask biometric fatigue indices, rendering manual listening prone to error.")
    ]
    for title, text in bullets_c2:
        p = tf_c2.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(13)
        r1.font.color.rgb = c_white
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(13)
        r2.font.color.rgb = c_gray

    # ==========================================
    # SLIDE 3: THE SOLUTION (3 columns)
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide3)
    add_header(slide3, "The Solution: Fusing Driver State with Telemetry")
    
    col_w = Inches(3.7)
    gap = Inches(0.4)
    top_pos = Inches(1.8)
    h_pos = Inches(4.8)
    
    steps = [
        ("01", "REAL-TIME SPEECH & SENTIMENT", c_red, [
            ("Acoustic Capture:", " Ingests raw driver radio streams, dynamically downmixing to 16kHz mono."),
            ("ASR Transcription:", " Runs custom Whisper models to transcribe radio clips in seconds."),
            ("Biometric Extraction:", " Computes speech volume (RMS) and silence pauses to detect exhaustion.")
        ]),
        ("02", "MULTIMODAL FUSION & CLASSIFICATION", c_cyan, [
            ("Dual-Sentiment AI:", " Fuses Audio emotion classification with Text sentiment classification."),
            ("Stress Formulas:", " Combines pace, volume, and keyword triggers to estimate stress (0-100)."),
            ("State Classification:", " Auto-labels drivers as CALM, STRESSED, or TIRED dynamically.")
        ]),
        ("03", "TELEMETRY CORRELATION HUD", c_white, [
            ("Live Ingestion:", " Integrates lap telemetry and sector split logs via Pandas."),
            ("Performance Delta:", " Compares post-stress performance against a strict calm-lap baseline."),
            ("Action Directives:", " Displays clear recommendations (Plan B pit swap, minimize radio traffic).")
        ])
    ]
    
    for i, (num, name, num_color, points) in enumerate(steps):
        left_pos = Inches(0.6) + i * (col_w + gap)
        create_card(slide3, left_pos, top_pos, col_w, h_pos, c_card)
        
        tb = slide3.shapes.add_textbox(left_pos + Inches(0.25), top_pos + Inches(0.25), col_w - Inches(0.5), h_pos - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Step Number
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.name = "Arial"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = num_color
        p_num.space_after = Pt(5)
        
        # Step Name
        p_name = tf.add_paragraph()
        p_name.text = name
        p_name.font.name = "Arial"
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = c_white
        p_name.space_after = Pt(20)
        
        for p_title, p_desc in points:
            p_pt = tf.add_paragraph()
            p_pt.space_after = Pt(10)
            r1 = p_pt.add_run()
            r1.text = p_title
            r1.font.bold = True
            r1.font.size = Pt(11)
            r1.font.color.rgb = c_white
            
            r2 = p_pt.add_run()
            r2.text = p_desc
            r2.font.size = Pt(11)
            r2.font.color.rgb = c_gray

    # ==========================================
    # SLIDE 4: TECH APPROACH - ML PIPELINE
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide4)
    add_header(slide4, "Tech Approach: Multimodal Sentiment Pipeline")
    
    # Left: Big Pipeline Diagram/Explanation
    create_card(slide4, Inches(0.6), Inches(1.8), Inches(6.5), Inches(4.8), c_card)
    tb_l = slide4.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.9), Inches(4.4))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    p = tf_l.paragraphs[0]
    p.text = "MULTIMODAL EMOTION PIPELINE"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(15)
    
    models = [
        ("ASR Transcription:", " openai/whisper-tiny", " (Fine-tuned on F1 radio corpus, fast CPU inference)."),
        ("Text Emotion:", " j-hartmann/emotion-english-distilroberta-base", " (Maps anger, fear, disgust, sadness)."),
        ("Audio Emotion (SER):", " superb/wav2vec2-base-superb-er", " (Classifies voice waveforms into angry/neutral/sad)."),
        ("Acoustic Stats:", " Waveform RMS & Silence ratios", " (Extracts speech volume & pauses/hesitations directly).")
    ]
    for m_name, m_model, m_desc in models:
        p = tf_l.add_paragraph()
        p.space_after = Pt(8)
        r1 = p.add_run()
        r1.text = m_name
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = c_white
        
        r2 = p.add_run()
        r2.text = m_model
        r2.font.bold = True
        r2.font.size = Pt(12)
        r2.font.color.rgb = c_cyan
        
        r3 = p.add_run()
        r3.text = m_desc
        r3.font.size = Pt(12)
        r3.font.color.rgb = c_gray

    # Right: Formula & Classification
    create_card(slide4, Inches(7.5), Inches(1.8), Inches(5.2), Inches(4.8), c_card)
    tb_r = slide4.shapes.add_textbox(Inches(7.8), Inches(2.0), Inches(4.6), Inches(4.4))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    p = tf_r.paragraphs[0]
    p.text = "THE FUSION ENGINE"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(25)
    
    # Formula Box in center
    p_f = tf_r.add_paragraph()
    p_f.text = "FUSION FORMULA:"
    p_f.font.bold = True
    p_f.font.size = Pt(12)
    p_f.font.color.rgb = c_white
    p_f.space_after = Pt(5)
    
    p_eq = tf_r.add_paragraph()
    p_eq.text = "Final Stress = 45% Voice + 30% Text + 15% Urgency + 10% Fatigue"
    p_eq.font.bold = True
    p_eq.font.size = Pt(12)
    p_eq.font.color.rgb = c_cyan
    p_eq.space_after = Pt(25)
    
    p_class = tf_r.add_paragraph()
    p_class.text = "DRIVER STATE DECISION MATRIX:"
    p_class.font.bold = True
    p_class.font.size = Pt(12)
    p_class.font.color.rgb = c_white
    p_class.space_after = Pt(10)
    
    states = [
        ("CALM:", " Low stress, low fatigue.", c_green),
        ("STRESSED:", " Stress score > 50.0.", c_red),
        ("TIRED:", " Fatigue score >= 50.0.", c_orange)
    ]
    for s_name, s_desc, s_color in states:
        p_s = tf_r.add_paragraph()
        p_s.space_after = Pt(8)
        r1 = p_s.add_run()
        r1.text = " •  " + s_name + " "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = s_color
        
        r2 = p_s.add_run()
        r2.text = s_desc
        r2.font.size = Pt(12)
        r2.font.color.rgb = c_gray

    # ==========================================
    # SLIDE 5: TECH APPROACH - CORRELATION
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide5)
    add_header(slide5, "Tech Approach: Telemetry & Correlation Engine")
    
    col_w2 = Inches(3.7)
    gap2 = Inches(0.4)
    top_pos2 = Inches(1.8)
    h_pos2 = Inches(4.8)
    
    telemetry_steps = [
        ("01", "INGESTION & CLEANING", c_white, [
            ("CSV Processing:", " Dynamic CSV ingestion parses lap-times and sector times using Pandas."),
            ("Deduplication:", " Wipes previous logs in SQL database for the active session to avoid overlaps."),
            ("Sector Backups:", " Computes fallback sector splits (Lap Time / 3) if telemetry lacks explicit sectors.")
        ]),
        ("02", "BASELINE CALCULATION", c_cyan, [
            ("Trigger Detection:", " Finds first transition lap where driver status shifts from Calm to Stressed/Tired."),
            ("Calm Baseline:", " Computes baseline by averaging last 3 calm laps immediately before stress trigger."),
            ("Anchor Point:", " Establishes a reference for subsequent pace degradation comparisons.")
        ]),
        ("03", "PACE LOSS CORRELATION", c_red, [
            ("Deterioration %:", " Formula: (Lap Time - Baseline) / Baseline * 100."),
            ("Sector Analytics:", " Pins performance loss to straight-line speed vs cornering sectors."),
            ("AI Directives:", " Renders mechanical wear hints (e.g. thermal tire degradation) & tactical actions.")
        ])
    ]
    
    for i, (num, name, num_color, points) in enumerate(telemetry_steps):
        left_pos = Inches(0.6) + i * (col_w2 + gap2)
        create_card(slide5, left_pos, top_pos2, col_w2, h_pos2, c_card)
        
        tb = slide5.shapes.add_textbox(left_pos + Inches(0.25), top_pos2 + Inches(0.25), col_w2 - Inches(0.5), h_pos2 - Inches(0.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Step Number
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.name = "Arial"
        p_num.font.size = Pt(28)
        p_num.font.bold = True
        p_num.font.color.rgb = num_color
        p_num.space_after = Pt(5)
        
        # Step Name
        p_name = tf.add_paragraph()
        p_name.text = name
        p_name.font.name = "Arial"
        p_name.font.size = Pt(13)
        p_name.font.bold = True
        p_name.font.color.rgb = c_white
        p_name.space_after = Pt(20)
        
        for p_title, p_desc in points:
            p_pt = tf.add_paragraph()
            p_pt.space_after = Pt(10)
            r1 = p_pt.add_run()
            r1.text = p_title
            r1.font.bold = True
            r1.font.size = Pt(11)
            r1.font.color.rgb = c_white
            
            r2 = p_pt.add_run()
            r2.text = p_desc
            r2.font.size = Pt(11)
            r2.font.color.rgb = c_gray

    # ==========================================
    # SLIDE 6: THE STAGE DEMO STORY
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide6)
    add_header(slide6, "Stage Demo: A 2-Minute Race Engineering Narrative")
    
    col_w3 = Inches(2.7)
    gap3 = Inches(0.25)
    top_pos3 = Inches(1.8)
    h_pos3 = Inches(4.8)
    
    demo_flow = [
        ("STEP 1", "ONE-CLICK START", c_white, [
            ("Deterministic Seed:", " Load Demo Session loads a pre-compiled 22-lap telemetry file."),
            ("Immediate Sync:", " Bypasses live model load latency on stage."),
            ("Starting State:", " Driver begins in CALM state, establishing telemetry baseline.")
        ]),
        ("STEP 2", "THE STRESS EVENT", c_red, [
            ("Lap 18 Radio:", " Select Lap 18 in transcripts log."),
            ("Voice Transcript:", " 'The tires are gone, I have no grip'"),
            ("Biometric Trigger:", " Fused stress score rises to 78%, flashing STRESSED state.")
        ]),
        ("STEP 3", "CORRELATED PACE LOSS", c_cyan, [
            ("Baseline Delta:", " Lap times deteriorate to 92.16s, triggering a +2.4% pace loss delta."),
            ("Track Diagnostics:", " Sector 2 (corners) shows a +2.50% degradation."),
            ("Biometric overlay:", " Stress HUD correlates to pace curve.")
        ]),
        ("STEP 4", "ENGINEER ACTION", c_green, [
            ("Tactical Directives:", " Race Engineer screen syncs alerts instantly."),
            ("Direct Actions:", " Recommends pit Plan B swap and radio silencing."),
            ("Conclusion:", " Real-time objective biometrics prevents pace bleed.")
        ])
    ]
    
    for i, (num, name, num_color, points) in enumerate(demo_flow):
        left_pos = Inches(0.6) + i * (col_w3 + gap3)
        create_card(slide6, left_pos, top_pos3, col_w3, h_pos3, c_card)
        
        tb = slide6.shapes.add_textbox(left_pos + Inches(0.15), top_pos3 + Inches(0.2), col_w3 - Inches(0.3), h_pos3 - Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        
        # Step Number
        p_num = tf.paragraphs[0]
        p_num.text = num
        p_num.font.name = "Arial"
        p_num.font.size = Pt(20)
        p_num.font.bold = True
        p_num.font.color.rgb = num_color
        p_num.space_after = Pt(5)
        
        # Step Name
        p_name = tf.add_paragraph()
        p_name.text = name
        p_name.font.name = "Arial"
        p_name.font.size = Pt(11)
        p_name.font.bold = True
        p_name.font.color.rgb = c_white
        p_name.space_after = Pt(15)
        
        for p_title, p_desc in points:
            p_pt = tf.add_paragraph()
            p_pt.space_after = Pt(8)
            r1 = p_pt.add_run()
            r1.text = p_title
            r1.font.bold = True
            r1.font.size = Pt(10)
            r1.font.color.rgb = c_white
            
            r2 = p_pt.add_run()
            r2.text = p_desc
            r2.font.size = Pt(10)
            r2.font.color.rgb = c_gray

    # ==========================================
    # SLIDE 7: TECH STACK & SYSTEM HIGHLIGHTS
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_background(slide7)
    add_header(slide7, "Tech Stack & Implementation Highlights")
    
    # Left Card: Technology Stack
    create_card(slide7, Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8), c_card)
    tb_ts = slide7.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_ts = tb_ts.text_frame
    tf_ts.word_wrap = True
    tf_ts.margin_left = tf_ts.margin_right = tf_ts.margin_top = tf_ts.margin_bottom = 0
    
    p = tf_ts.paragraphs[0]
    p.text = "PRODUCTION-READY TECH STACK"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(20)
    
    stack = [
        ("Frontend Application:", " Next.js 15+, React 19, TypeScript, Tailwind CSS v4, Recharts (responsive double-axis HUD charts)."),
        ("Backend Services:", " Python 3.14, FastAPI, SQLAlchemy, SQLite, Uvicorn."),
        ("ML Engineering:", " PyTorch, Hugging Face Hub (Whisper-Tiny, DistilRoBERTa, Wav2Vec2-ER)."),
        ("Data Analysis:", " Pandas, NumPy (biometric & CSV telemetry alignment).")
    ]
    for title, text in stack:
        p = tf_ts.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = c_white
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(12)
        r2.font.color.rgb = c_gray

    # Right Card: Hackathon Hardening
    create_card(slide7, Inches(6.9), Inches(1.8), Inches(5.8), Inches(4.8), c_card)
    tb_hh = slide7.shapes.add_textbox(Inches(7.2), Inches(2.0), Inches(5.2), Inches(4.4))
    tf_hh = tb_hh.text_frame
    tf_hh.word_wrap = True
    tf_hh.margin_left = tf_hh.margin_right = tf_hh.margin_top = tf_hh.margin_bottom = 0
    
    p = tf_hh.paragraphs[0]
    p.text = "HACKATHON-READY HARDENING"
    p.font.name = "Arial"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = c_red
    p.space_after = Pt(20)
    
    hardening = [
        ("100% Offline Edge Execution:", " All transcriptions, SER voice evaluations, and telemetry calibrations run fully locally. Zero external API risk."),
        ("Dynamic React Context:", " Shared SessionContext matches updates across all 4 dashboards in real-time."),
        ("Fast CPU Inference Singletons:", " Warmup caches ensure subsequent inferences resolve in under 5 seconds."),
        ("Deterministic Demonstration:", " Load Demo Session sews telemetry + transcriptions instantly for flawless live pitches.")
    ]
    for title, text in hardening:
        p = tf_hh.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = title
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = c_white
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(12)
        r2.font.color.rgb = c_gray

    # Save Presentation
    output_filename = "PitPulse_Hackathon.pptx"
    prs.save(output_filename)
    print(f"Presentation saved successfully as '{output_filename}'!")

if __name__ == "__main__":
    create_deck()
